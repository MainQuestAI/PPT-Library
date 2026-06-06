#!/usr/bin/env python3
from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError as exc:  # pragma: no cover - exercised by users without demo extra.
    raise SystemExit(
        "python-pptx is required for demo deck generation. "
        "Run: uv run --extra demo python scripts/create_demo_decks.py --output /tmp/ppt-lib-demo-decks"
    ) from exc


@dataclass(frozen=True)
class DemoDeck:
    filename: str
    title: str
    subtitle: str
    accent: tuple[int, int, int]
    slides: list[tuple[str, str, str]]


DEMO_DECKS: list[DemoDeck] = [
    DemoDeck(
        filename="synthetic-retail-growth-playbook.pptx",
        title="Retail Growth Playbook",
        subtitle="Synthetic demo deck for reusable sales content",
        accent=(245, 158, 11),
        slides=[
            ("现状挑战", "门店会员运营分散，活动复盘难以沉淀为可复用方案。", "problem"),
            ("解决方案", "建立会员分层、活动编排、复购触达的一体化运营能力。", "solution"),
            ("业务架构图", "数据采集 -> 人群洞察 -> 触达策略 -> ROI 复盘。", "architecture"),
            ("ROI 价值", "复用高胜率页面，缩短方案准备时间，提升提案一致性。", "roi"),
            ("下一步共创", "用三场工作坊确认样例库、标签口径和复用追踪方式。", "cta"),
        ],
    ),
    DemoDeck(
        filename="synthetic-manufacturing-service-blueprint.pptx",
        title="Manufacturing Service Blueprint",
        subtitle="Synthetic demo deck for service transformation proposals",
        accent=(37, 99, 235),
        slides=[
            ("痛点", "设备服务流程长，售前方案难以快速说明端到端价值。", "problem"),
            ("方案蓝图", "工单接入、知识推荐、现场协同、客户反馈闭环。", "solution"),
            ("流程图", "服务请求 -> 诊断 -> 派工 -> 备件 -> 复盘。", "architecture"),
            ("案例实践", "某合成制造企业用标准化页面讲清楚服务升级路径。", "case"),
            ("收益", "减少重复造页，沉淀可复用行业表达。", "roi"),
        ],
    ),
    DemoDeck(
        filename="synthetic-public-sector-data-platform.pptx",
        title="Public Sector Data Platform",
        subtitle="Synthetic demo deck for local asset intelligence workflows",
        accent=(16, 185, 129),
        slides=[
            ("挑战", "多部门材料版本分散，政策汇报页面复用效率低。", "problem"),
            ("建设方案", "资料入库、关键页识别、人工复核、搜索复用。", "solution"),
            ("体系架构", "本地索引库、标签管理、Agent 调用、审查导出。", "architecture"),
            ("价值证明", "关键页、战绩、搜索排序形成资产经营闭环。", "roi"),
            ("行动计划", "先用合成样例验证，再替换为用户授权资料。", "cta"),
        ],
    ),
]


def main() -> int:
    parser = argparse.ArgumentParser(description="Create synthetic PPTX decks for PPT Library demos.")
    parser.add_argument("--output", type=Path, required=True, help="Directory where synthetic PPTX files will be written.")
    args = parser.parse_args()

    output_dir = args.output.expanduser().resolve()
    output_dir.mkdir(parents=True, exist_ok=True)

    written: list[Path] = []
    for deck in DEMO_DECKS:
        path = output_dir / deck.filename
        write_demo_deck(path, deck)
        written.append(path)

    for path in written:
        print(path)
    return 0


def write_demo_deck(path: Path, deck: DemoDeck) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]

    add_cover_slide(presentation, blank_layout, deck)
    for index, (title, body, role) in enumerate(deck.slides, start=2):
        add_content_slide(presentation, blank_layout, deck, index, title, body, role)

    presentation.save(path)


def add_cover_slide(presentation: Presentation, blank_layout, deck: DemoDeck) -> None:
    slide = presentation.slides.add_slide(blank_layout)
    accent = RGBColor(*deck.accent)
    add_background(slide, accent)

    title = slide.shapes.add_textbox(Inches(0.85), Inches(1.35), Inches(8.6), Inches(0.8))
    title_frame = title.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = deck.title
    title_run.font.size = Pt(34)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(17, 24, 39)

    subtitle = slide.shapes.add_textbox(Inches(0.88), Inches(2.32), Inches(7.2), Inches(0.45))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.clear()
    subtitle_run = subtitle_frame.paragraphs[0].add_run()
    subtitle_run.text = deck.subtitle
    subtitle_run.font.size = Pt(15)
    subtitle_run.font.color.rgb = RGBColor(75, 85, 99)

    badge = slide.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(3.4), Inches(0.55))
    badge_frame = badge.text_frame
    badge_frame.clear()
    badge_p = badge_frame.paragraphs[0]
    badge_p.alignment = PP_ALIGN.CENTER
    badge_run = badge_p.add_run()
    badge_run.text = "Synthetic demo content"
    badge_run.font.size = Pt(13)
    badge_run.font.bold = True
    badge_run.font.color.rgb = RGBColor(255, 255, 255)
    badge.fill.solid()
    badge.fill.fore_color.rgb = accent
    badge.line.color.rgb = accent

    add_signal_card(slide, deck, "Key pages", "Searchable, reviewable, reusable", Inches(8.3), Inches(1.35))
    add_signal_card(slide, deck, "Deal signals", "Usage and win/loss context", Inches(8.3), Inches(3.0))
    add_signal_card(slide, deck, "Agent-ready", "Stable JSON and local assets", Inches(8.3), Inches(4.65))


def add_content_slide(
    presentation: Presentation,
    blank_layout,
    deck: DemoDeck,
    page_number: int,
    title_text: str,
    body_text: str,
    role: str,
) -> None:
    slide = presentation.slides.add_slide(blank_layout)
    accent = RGBColor(*deck.accent)
    add_background(slide, accent)

    role_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.55), Inches(1.6), Inches(0.36))
    role_frame = role_box.text_frame
    role_frame.clear()
    role_p = role_frame.paragraphs[0]
    role_p.alignment = PP_ALIGN.CENTER
    role_run = role_p.add_run()
    role_run.text = role
    role_run.font.size = Pt(11)
    role_run.font.bold = True
    role_run.font.color.rgb = RGBColor(255, 255, 255)
    role_box.fill.solid()
    role_box.fill.fore_color.rgb = accent
    role_box.line.color.rgb = accent

    title = slide.shapes.add_textbox(Inches(0.75), Inches(1.12), Inches(7.8), Inches(0.75))
    title_frame = title.text_frame
    title_frame.clear()
    title_run = title_frame.paragraphs[0].add_run()
    title_run.text = title_text
    title_run.font.size = Pt(30)
    title_run.font.bold = True
    title_run.font.color.rgb = RGBColor(17, 24, 39)

    body = slide.shapes.add_textbox(Inches(0.78), Inches(2.12), Inches(6.8), Inches(1.15))
    body_frame = body.text_frame
    body_frame.word_wrap = True
    body_frame.clear()
    body_run = body_frame.paragraphs[0].add_run()
    body_run.text = body_text
    body_run.font.size = Pt(19)
    body_run.font.color.rgb = RGBColor(31, 41, 55)

    add_visual_stack(slide, deck, role)
    add_footer(slide, deck.title, page_number)


def add_background(slide, accent: RGBColor) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(249, 250, 251)

    ribbon = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = accent
    ribbon.line.color.rgb = accent


def add_signal_card(slide, deck: DemoDeck, title: str, body: str, left, top) -> None:
    accent = RGBColor(*deck.accent)
    card = slide.shapes.add_shape(1, left, top, Inches(3.6), Inches(1.05))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(229, 231, 235)

    text = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.18), Inches(3.1), Inches(0.62))
    frame = text.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = accent
    p2 = frame.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(75, 85, 99)


def add_visual_stack(slide, deck: DemoDeck, role: str) -> None:
    accent = RGBColor(*deck.accent)
    labels = {
        "problem": ["Fragmented source", "Manual reuse", "Slow proposal"],
        "solution": ["Tag", "Review", "Reuse"],
        "architecture": ["Local files", "Index", "Agent JSON"],
        "case": ["Scenario", "Slide", "Outcome"],
        "roi": ["Reuse count", "Won count", "Win rate"],
        "cta": ["Confirm", "Index", "Improve"],
    }.get(role, ["Source", "Signal", "Search"])
    for index, label in enumerate(labels):
        left = Inches(8.25)
        top = Inches(1.35 + index * 1.35)
        box = slide.shapes.add_shape(1, left, top, Inches(3.75), Inches(0.85))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 255, 255)
        box.line.color.rgb = accent
        text = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.22), Inches(3.25), Inches(0.36))
        frame = text.text_frame
        frame.clear()
        run = frame.paragraphs[0].add_run()
        run.text = label
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(31, 41, 55)


def add_footer(slide, deck_title: str, page_number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.78), Inches(6.85), Inches(6.0), Inches(0.28))
    frame = footer.text_frame
    frame.clear()
    run = frame.paragraphs[0].add_run()
    run.text = f"{deck_title} · synthetic page {page_number}"
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(107, 114, 128)


if __name__ == "__main__":
    raise SystemExit(main())
